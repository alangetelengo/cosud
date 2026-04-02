# -*- coding: utf-8 -*-
"""
Script pour mettre à jour le contenu de la présentation GED
en conservant la charte graphique ACSI (couleurs, design).
"""
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER

# Chemins
SOURCE_PPTX = r"D:\GED\PRESENTATION  GED - SI.pptx"
OUTPUT_PPTX = r"D:\GED\PRESENTATION  GED - SI.pptx"  # Fichier modifié (sauvegarde in-place)
PLAN_MD = Path(__file__).parent.parent / "docs" / "PLAN_PRESENTATION_GED.md"

# Contenu des slides (ordre du plan)
SLIDES_CONTENT = [
    # Slide 1 - Page de titre
    {
        "title": "GED",
        "body": "Gestion Électronique des Documents\n\nPropulsé par l'ACSI\n\nPrésenté par : [Votre nom] – Ingénieur SI\nDate : [Date de présentation]"
    },
    # Slide 2 - Contexte et objectif
    {
        "title": "Contexte et objectif",
        "body": "Problématique\n• Multiplication des documents (papier et numérique)\n• Difficultés de classement et de recherche\n• Risque de perte ou de duplication\n• Validation hiérarchique peu tracée\n• Partage d'informations non structuré\n\nObjectif d'GED\nCentraliser, sécuriser et organiser les documents professionnels avec :\n• Un dépôt simple et structuré\n• Une recherche rapide\n• Une chaîne de validation traçable\n• Des droits d'accès maîtrisés"
    },
    # Slide 3 - Vue d'ensemble
    {
        "title": "Vue d'ensemble du système",
        "body": "GED en quelques points\n• Type : Application web GED (SaaS)\n• Technologie : Laravel, PHP, Tailwind CSS\n• Accès : Navigateur web, authentification sécurisée\n• Données : Hébergées en interne, contrôle total\n• Évolutivité : Architecture modulaire\n\nBénéfices principaux\n✓ Réduction du papier\n✓ Traçabilité des actions\n✓ Accès contrôlé par rôles\n✓ Workflow de validation hiérarchique\n✓ Notifications automatiques (email, in-app, SMS)"
    },
    # Slide 4 - Architecture fonctionnelle
    {
        "title": "Architecture fonctionnelle",
        "body": "Modules principaux\n\n• Tableau de bord : Vue adaptée DG / Structure / User\n• Documents : Dépôt, consultation, workflow\n• Dossiers : Arborescence, partages, favoris\n• Recherche : Full-text (titre, ref…)\n• Corbeille : Restauration documents supprimés\n• Notifications : Alertes dépôt, validation, rejet\n\nAdministration (réservé aux admins)\n• Utilisateurs | Structures | Rôles/Permissions\n• Plan de classement | Types | Workflow | Audit"
    },
    # Slide 5 - Tableaux de bord
    {
        "title": "Tableaux de bord adaptés",
        "body": "Trois vues selon le profil\n\n• DG / Admin : Stats org., documents par structure, en attente, récents\n• Responsable de structure : Docs de sa direction, en attente, récents\n• Utilisateur : Mes documents, favoris, derniers dossiers\n\nIndicateurs clés\n• Nombre de documents (brouillon, en attente, validés, rejetés, archivés)\n• Répartition par structure\n• Derniers dépôts\n• Documents en attente de validation"
    },
    # Slide 6 - Gestion documents (1/2)
    {
        "title": "Gestion des documents (1/2)",
        "body": "Cycle de vie\nDÉPÔT → BROUILLON → ENVOI VALIDATION → EN ATTENTE → VALIDÉ / REJETÉ → ARCHIVÉ\n\nDépôt de document\n• Glisser-déposer ou sélection de fichier\n• Formats : PDF, Word, Excel, images\n• Métadonnées : titre, référence, mots-clés, description\n• Extraction automatique des métadonnées PDF\n• Suggestion du type selon l'extension\n• Classement dans un dossier (optionnel)\n\nGestion des versions : Historique et traçabilité"
    },
    # Slide 7 - Gestion documents (2/2)
    {
        "title": "Gestion des documents (2/2)",
        "body": "Workflow de validation\n• Validation directe (admin/DG) ou workflow par étapes\n• Validation hiérarchique : structure du créateur jusqu'au DG\n• Approuver ou rejeter avec motif obligatoire\n• Notifications aux validateurs et au créateur\n\nStatuts\n• Brouillon | En attente | Validé | Rejeté | Archivé"
    },
    # Slide 8 - Dossiers
    {
        "title": "Dossiers et plan de classement",
        "body": "Arborescence des dossiers\n• Plan de classement paramétrable par l'admin\n• Hiérarchie illimitée (ex. : Finance > Comptabilité > Factures clients)\n• Liaison dossiers ↔ structures (Direction responsable)\n• Dossiers confidentiels et alertes SMS\n\nPartages\n• Partage par utilisateur avec droits (lecture, écriture, suppression)\n• Date d'expiration optionnelle\n• Propriétaire = responsable de la structure (acteur métier)\n\nFavoris : Accès rapide aux dossiers les plus utilisés"
    },
    # Slide 9 - Sécurité
    {
        "title": "Sécurité et traçabilité",
        "body": "Gestion des accès\n• Rôles : Admin, DG, Utilisateur\n• Permissions granulaires (documents, dossiers, types, utilisateurs…)\n• Structure hiérarchique (chaque direction a son périmètre)\n• Dossiers confidentiels : permission spéciale\n\nJournal d'audit\n• Enregistrement des actions sensibles\n• Identification de l'acteur et de l'heure\n\nIntégrité : Empreinte SHA-256 des fichiers"
    },
    # Slide 10 - Notifications
    {
        "title": "Notifications",
        "body": "Canaux\n• Email : Dépôt, validation, rejet\n• In-app : Notification dans l'interface\n• SMS : Dossiers importants/confidentiels (optionnel, Vonage)\n\nActeurs notifiés (acteurs métier uniquement)\n• Propriétaire du dossier\n• Utilisateurs en partage lecture\n• Validateurs lors de l'envoi en validation\n• Créateur lors d'un approuvé/rejet"
    },
    # Slide 11 - Administration
    {
        "title": "Administration et paramétrage",
        "body": "Paramètres (réservés aux admins)\n• Structures : organigramme, responsables\n• Rôles et permissions : matrice des droits\n• Plan de classement : arborescence des dossiers\n• Types de documents : PDF, Word, Excel… (extensions, taille max)\n• Types de dossiers : Administration, Finance, Projet…\n• Types de métadonnées : champs indexables\n• Workflow : étapes de validation par type de document\n• Journal d'audit : consultation des actions"
    },
    # Slide 12 - Technologies
    {
        "title": "Technologies et déploiement",
        "body": "Stack technique\n• Backend : Laravel (PHP)\n• Frontend : Blade, Tailwind CSS, Alpine.js\n• Base de données : MySQL / MariaDB\n• Stockage : Système de fichiers (local ou S3-compatible)\n\nDéploiement\n• Hébergement sur serveur interne\n• SSL recommandé pour la production\n• Sauvegardes régulières des données et des fichiers"
    },
    # Slide 13 - Prochaines étapes
    {
        "title": "Prochaines étapes",
        "body": "• Formation des utilisateurs clés\n• Migration des documents existants (si applicable)\n• Réglage des workflows selon les processus métier\n• Mise en production progressive par direction"
    },
    # Slide 14 - Conclusion
    {
        "title": "Conclusion",
        "body": "GED est une solution GED interne, alignée sur l'organisation de l'ACSI :\n• Gestion des documents centralisée\n• Validation hiérarchique traçable\n• Sécurité et audit intégrés\n• Interface moderne et responsive\n\nContact : [Votre email]"
    },
]


def set_shape_text(shape, text):
    """Remplace le texte d'une forme en conservant le style existant."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0] if tf.paragraphs else tf.add_paragraph()
    p.text = text


def update_slide_content(slide, content):
    """Met à jour le contenu d'une slide en préservant le design (couleurs, polices)."""
    title = content.get("title", "")
    body = content.get("body", "")
    
    title_done = False
    body_done = False
    
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        # Identifier titre vs corps via le type de placeholder
        if shape.is_placeholder:
            try:
                ph_type = shape.placeholder_format.type
                if ph_type == PP_PLACEHOLDER.TITLE:
                    set_shape_text(shape, title)
                    title_done = True
                elif ph_type == PP_PLACEHOLDER.BODY:
                    if not title_done:
                        set_shape_text(shape, title)
                        title_done = True
                    else:
                        set_shape_text(shape, body)
                        body_done = True
                elif ph_type == PP_PLACEHOLDER.SUBTITLE:
                    set_shape_text(shape, body)
                    body_done = True
            except (AttributeError, KeyError):
                pass
        if title_done and body_done:
            break
    
    # Fallback : si pas de placeholders reconnus, utiliser l'ordre (1ère = titre, 2e = corps)
    if not title_done or not body_done:
        text_shapes = [s for s in slide.shapes if s.has_text_frame]
        for i, shp in enumerate(text_shapes):
            if i == 0 and not title_done:
                set_shape_text(shp, title)
                title_done = True
            elif i >= 1 and not body_done:
                set_shape_text(shp, body)
                body_done = True
                break


def main():
    if not os.path.exists(SOURCE_PPTX):
        print(f"ERREUR : Fichier non trouvé : {SOURCE_PPTX}")
        return 1
    
    # Sauvegarde de l'original avant modification
    backup_path = SOURCE_PPTX.replace(".pptx", "_backup_avant_modif.pptx")
    if not os.path.exists(backup_path):
        import shutil
        shutil.copy2(SOURCE_PPTX, backup_path)
        print(f"Sauvegarde créée : {backup_path}")
    
    print(f"Ouverture de {SOURCE_PPTX}...")
    prs = Presentation(SOURCE_PPTX)
    
    output = OUTPUT_PPTX
    
    nb_existing = len(prs.slides)
    nb_content = len(SLIDES_CONTENT)
    
    print(f"Slides existantes : {nb_existing}")
    print(f"Slides à appliquer : {nb_content}")
    
    # Stratégie : remplacer le contenu des slides existantes
    # Si plus de slides existantes que de contenu, on met à jour les premières
    # Si moins, on ajoute des slides avec le layout de la dernière
    
    for i, content in enumerate(SLIDES_CONTENT):
        if i < nb_existing:
            # Remplacer le contenu de la slide existante
            slide = prs.slides[i]
            print(f"  Mise à jour slide {i+1} : {content['title'][:40]}...")
            update_slide_content(slide, content)
        else:
            # Ajouter une nouvelle slide avec le même layout que la dernière
            last_slide = prs.slides[-1]
            slide_layout = last_slide.slide_layout
            new_slide = prs.slides.add_slide(slide_layout)
            # Ajouter titre et corps
            for shape in new_slide.shapes:
                if shape.has_text_frame:
                    if shape.placeholder_format and shape.placeholder_format.idx == 0:
                        set_shape_text(shape, content["title"])
                    else:
                        set_shape_text(shape, content["body"])
                        break
            print(f"  Ajout slide {i+1} : {content['title'][:40]}...")
    
    prs.save(output)
    print(f"\nPrésentation sauvegardée : {output}")
    return 0


if __name__ == "__main__":
    exit(main())
