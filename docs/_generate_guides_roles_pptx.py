"""
Génère un guide PowerPoint permanent par rôle COSUD.
Sortie : docs/guides-roles/Guide-<role>.pptx
Aucun nom de personne — uniquement le libellé du rôle.
"""

from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

DOCS = Path(__file__).resolve().parent
OUT_DIR = DOCS / "guides-roles"

GREEN = RGBColor(0x06, 0xA2, 0x69)
DARK = RGBColor(0x0F, 0x17, 0x2A)
GRAY = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)


def _set_run(run, *, size=18, bold=False, color=DARK, italic=False):
    run.font.name = "Calibri"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_bg(slide, color: RGBColor):
    fill = slide.shapes.add_shape(
        1,  # rectangle
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    fill.fill.solid()
    fill.fill.fore_color.rgb = color
    fill.line.fill.background()
    # send to back
    spTree = slide.shapes._spTree
    sp = fill._element
    spTree.remove(sp)
    spTree.insert(2, sp)


def _banner(slide, title: str, subtitle: str | None = None):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    shape.fill.solid()
    shape.fill.fore_color.rgb = GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    _set_run(run, size=28, bold=True, color=WHITE)
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.margin_top = Inches(0.2)
    if subtitle:
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = subtitle
        _set_run(r2, size=14, color=WHITE)


def _bullets(slide, items: list[str], top=1.4, left=0.5, width=12.3, size=16):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(8)
        run = p.add_run()
        run.text = "•  " + text
        _set_run(run, size=size, color=DARK)


def _footer(slide, role_label: str):
    box = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.35))
    tf = box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = f"COSUD — ACSI | Guide permanent — {role_label}"
    _set_run(run, size=10, italic=True, color=GRAY)


def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, role_label: str, mission: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(3))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "COSUD"
    _set_run(r, size=20, bold=True, color=GREEN)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"Guide utilisateur — {role_label}"
    _set_run(r2, size=36, bold=True, color=WHITE)
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = mission
    _set_run(r3, size=18, color=LIGHT)
    p3.space_before = Pt(18)
    p4 = tf.add_paragraph()
    r4 = p4.add_run()
    r4.text = "Guide permanent de l’application (droits et parcours métier)"
    _set_run(r4, size=14, italic=True, color=GRAY)
    p4.space_before = Pt(24)


def add_section(prs, role_label: str, title: str, bullets: list[str]):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, WHITE)
    _banner(slide, title, role_label)
    _bullets(slide, bullets)
    _footer(slide, role_label)


# ---------------------------------------------------------------------------
# Contenu par rôle (aucune personne nommée)
# ---------------------------------------------------------------------------

GUIDES: list[dict] = [
    {
        "file": "Guide-secretaire-direction.pptx",
        "label": "Secrétaire de direction",
        "mission": "Enregistrer et suivre les courriers d’arrivée et de départ pour la direction.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Tableau de bord, Documents, Dossiers, Courriers, Registre Arrivée, Registre Départ.",
                    "Voit les factures et les MAD (permissions voir-factures + voir-depenses).",
                    "Peut créer, modifier, transmettre, archiver et réceptionner des courriers.",
                    "Condition : être rattaché(e) au secrétariat (structure SEC-…) pour ouvrir « Nouveau courrier ».",
                ],
            ),
            (
                "Tâche principale — Courrier arrivée",
                [
                    "Courriers → Nouveau courrier arrivée.",
                    "Saisir le n° de registre (papier), l’objet, les dates, le type, puis joindre le scan (PDF/JPG obligatoire).",
                    "Choisir le type adapté : facture, MAD, demande, administratif, etc.",
                    "Enregistrer : le circuit démarre automatiquement selon le type.",
                    "Contrôler dans le Registre Arrivée que la ligne apparaît avec le bon n°.",
                ],
            ),
            (
                "Selon le type",
                [
                    "Facture : fournisseur dans le référentiel, montant, téléphone, service demandeur — obligatoires.",
                    "MAD : service demandeur obligatoire ; pas de montant ni de fiche fournisseur.",
                    "Demande : téléphone de l’expéditeur obligatoire (notifications).",
                    "Autres types : objet + scan + n° registre suffisent en général.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Un n° de registre déjà utilisé est refusé (doublon).",
                    "Sans scan, l’enregistrement arrivée est bloqué.",
                    "Après enregistrement, COSUD ouvre la fiche du courrier : revenir à « Nouveau » pour enchaîner.",
                    "Ne pas confondre avec « Reprise des factures prestataires » (module hors circuit / historique).",
                ],
            ),
        ],
    },
    {
        "file": "Guide-responsable-dossiers-prestataires.pptx",
        "label": "Responsable dossiers prestataires",
        "mission": "Piloter les factures fournisseurs / prestataires, le référentiel et le suivi des dettes.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus clés : Courriers (factures), Suivi de factures, Fournisseurs ou prestataires, Reprise des factures, Moratoires (consultation).",
                    "Voit les factures ; ne voit pas les MAD en liste (pas de permission voir-depenses).",
                    "Peut créer des fiches fournisseur / prestataire / partenaire et des factures de régularisation.",
                    "Peut créer des courriers (si structure secrétariat) — typiquement les factures.",
                ],
            ),
            (
                "Référentiel fournisseurs / prestataires",
                [
                    "Menu Fournisseurs ou prestataires → Créer.",
                    "Renseigner au minimum : nom, type (fournisseur / prestataire / partenaire).",
                    "Téléphone et e-mail facilitent les SMS et le préremplissage sur les factures.",
                    "Une facture ne peut pas être enregistrée sans fiche active dans ce référentiel.",
                    "Mettre à jour ou désactiver une fiche plutôt que créer des doublons de noms.",
                ],
            ),
            (
                "Factures dans le registre",
                [
                    "Nouveau courrier arrivée → type Facture prestataire.",
                    "Choisir le fournisseur dans la liste, saisir montant, service demandeur, téléphone, scan.",
                    "À l’enregistrement : circuit « facture prestataire » → attente Bon pour accord DG.",
                    "Suivre l’avancement sur la fiche courrier (étapes) et dans Suivi de factures.",
                ],
            ),
            (
                "Reprise des factures & dettes",
                [
                    "Module « Reprise des factures prestataires » : stock historique hors circuit (n° REG-…).",
                    "Ne pas l’utiliser pour une pièce du registre papier courant (utiliser Courriers).",
                    "Les montants alimentent les dettes fournisseurs (consultables via Suivi / Moratoires).",
                    "Moratoires : consultation uniquement ; la création relève du responsable suivi des dépenses.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Créer la fiche fournisseur avant la saisie de facture si le fournisseur est absent.",
                    "Classer la pièce dans le bon dossier prestataire lorsque la procédure le demande.",
                    "Vérifier les doublons de références facture fournisseur (COSUD peut bloquer).",
                ],
            ),
        ],
    },
    {
        "file": "Guide-responsable-suivi-depenses.pptx",
        "label": "Responsable suivi des dépenses",
        "mission": "Suivre les MAD, les paiements, les bordereaux et les plans de moratoire.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Courriers (MAD), Suivi de dépense, Bordereau de transmission, Moratoires, Reprise des factures (paiement).",
                    "Voit les MAD ; ne voit pas les factures en liste (pas de permission voir-factures).",
                    "Peut enregistrer des paiements (suivi) et payer une régularisation.",
                    "Seul rôle métier (hors DG/admin) habilité à créer et mettre à jour les moratoires.",
                ],
            ),
            (
                "MAD — enregistrement et suivi",
                [
                    "Nouveau courrier arrivée → type MAD (si habilité secrétariat).",
                    "Champs clés : n° registre, objet, service demandeur, scan.",
                    "Circuit identique aux factures (Bon pour accord DG, puis suite paiement).",
                    "Suivre dans Courriers / Suivi de dépense selon l’étape.",
                ],
            ),
            (
                "Suivi de dépense & paiements",
                [
                    "Menu Suivi de dépense : état des paiements liés aux dossiers.",
                    "Enregistrer un paiement lorsque la procédure et les pièces le permettent.",
                    "Bordereau de transmission : préparer / consulter les transmissions liées au flux.",
                    "Reprise des factures : action « payer » sur une régularisation éligible.",
                ],
            ),
            (
                "Moratoires",
                [
                    "Prérequis : factures saisies avec montants → dette fournisseur > 0.",
                    "Moratoires → Créer : choisir le fournisseur éligible, échéances, joindre l’instruction DG.",
                    "Ne créer un moratoire que sur instruction claire et dettes cohérentes.",
                    "Mettre à jour les échéances (paiement / statut) au fil de l’exécution du plan.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Les MAD n’alimentent pas la dette fournisseur (pas de montant facture).",
                    "Coordonner avec le responsable dossiers prestataires pour les montants facture.",
                    "Un moratoire actif empêche en général un second plan pour le même fournisseur.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-particulier-dg.pptx",
        "label": "Particulier DG",
        "mission": "Appuyer le cabinet DG sur le suivi des courriers, factures, dettes et notifications.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Courriers (factures + MAD), Suivi de dépense, Suivi de factures, Moratoires (vue), Bordereau.",
                    "Peut créer et gérer des courriers comme le secrétariat (si structure SEC).",
                    "Reçoit les notifications de circuit (cloche) pour suivre les dossiers en cours.",
                    "N’exécute pas à la place du DG les étapes « Bon pour accord » / signature chèque (rôle dg requis).",
                ],
            ),
            (
                "Travail quotidien",
                [
                    "Surveiller la cloche et la liste Courriers pour les dossiers à relayer au DG.",
                    "Vérifier la qualité des saisies (n°, scan, objet) avant présentation.",
                    "Consulter Suivi de factures et dettes pour préparer les arbitrages.",
                    "Utiliser les registres Arrivée / Départ pour les recherches rapides.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Distinguer notification (être informé) et action circuit (être l’acteur de l’étape).",
                    "Les dossiers confidentiels restreignent les notifications élargies.",
                    "Ne pas avancer une étape DG sans le compte DG.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-particulier-ac.pptx",
        "label": "Particulier AC",
        "mission": "Appuyer l’agent comptable sur le suivi des courriers et des transmissions.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Courriers (factures + MAD), Bordereau de transmission, Documents, Dossiers.",
                    "Peut créer / transmettre des courriers (secrétariat) si structure adaptée.",
                    "Reçoit les notifications d’étapes du circuit facture / MAD (suivi parallèle).",
                    "N’établit pas le chèque à la place de l’agent comptable (étape réservée au rôle AC).",
                ],
            ),
            (
                "Travail quotidien",
                [
                    "Suivre les dossiers en attente côté comptabilité via la cloche et Courriers.",
                    "Préparer les éléments pour l’AC (montants, fournisseur, instructions DG visibles sur la fiche).",
                    "Consulter / préparer les bordereaux de transmission.",
                    "Vérifier que les scans et références sont complets avant traitement AC.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "L’action « AC établit le chèque » doit être faite avec un compte agent comptable.",
                    "Rester aligné sur les instructions DG affichées sur la fiche courrier.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-agent-comptable.pptx",
        "label": "Agent comptable",
        "mission": "Établir les chèques et faire avancer le circuit de paiement des factures / MAD.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Courriers (factures + MAD), Bordereau de transmission, Documents, Dossiers.",
                    "Acteur clé de l’étape « AC établit le chèque → envoi DG » du circuit facture prestataire.",
                    "Peut créer des courriers si rattaché au secrétariat ; sinon travail principal = traitement circuit.",
                    "Pas d’accès dédié suivi-factures / moratoires create dans le seed standard (hors consultation via courriers).",
                ],
            ),
            (
                "Circuit facture / MAD — votre étape",
                [
                    "Après Bon pour accord DG : le dossier arrive à l’étape AC.",
                    "Ouvrir la fiche courrier → lire les instructions DG et le montant.",
                    "Établir le chèque hors COSUD si besoin, puis enregistrer l’action dans COSUD (envoi au DG pour signature).",
                    "Après signature DG : reprendre pour la décharge / preuve selon l’étape affichée.",
                ],
            ),
            (
                "Paiement reliquat (si activé sur le dossier)",
                [
                    "Sur une facture déjà partiellement traitée, un paiement de reliquat peut être proposé selon les règles métier.",
                    "Vérifier le reliquat et l’absence de blocage (ex. moratoire) avant d’enregistrer.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Ne pas signer à la place du DG (étape signature chèque = rôle dg).",
                    "Respecter le mouvement physique du dossier (attente arrivée) indiqué par le circuit.",
                    "Bordereau de transmission : utile pour tracer les envois.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-caissier.pptx",
        "label": "Caissier",
        "mission": "Participer au flux courriers / dépenses et aux transmissions selon les étapes assignées.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Courriers (factures + MAD), Bordereau de transmission, Documents, Dossiers.",
                    "Droits proches du secrétariat courrier + vue factures et MAD.",
                    "Intervient lorsque le circuit ou la procédure de caisse le sollicite.",
                ],
            ),
            (
                "Travail quotidien",
                [
                    "Consulter la cloche et la liste Courriers pour les dossiers à traiter.",
                    "Vérifier montants et pièces avant toute opération de caisse liée au dossier.",
                    "Utiliser le bordereau de transmission pour les envois documentés.",
                    "Ne pas avancer une étape réservée à un autre rôle (DG, AC) sans le bon compte.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Les paiements formalisés « suivi de dépense » relèvent surtout du responsable suivi des dépenses.",
                    "En cas de doute sur l’étape affichée, ouvrir la fiche et lire l’aide de l’étape / l’historique.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-dg.pptx",
        "label": "Directeur général",
        "mission": "Instruire, donner le Bon pour accord, signer et arbitrer l’ensemble des flux COSUD.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Accès large métier : Courriers, registres, suivi factures / dépenses, fournisseurs, régularisation, moratoires, utilisateurs.",
                    "Pas des menus d’administration technique réservés à l’administrateur (paramètres avancés, corbeille, etc.).",
                    "Acteur des étapes ACTEUR_DG : Bon pour accord, signature chèque, instructions.",
                ],
            ),
            (
                "Circuit facture / MAD",
                [
                    "Après enregistrement secrétariat : étape « Bon pour accord / instructions DG ».",
                    "Ouvrir la fiche → donner le Bon pour accord et saisir les instructions (montant, modalités…).",
                    "L’AC est alors notifié pour établir le chèque.",
                    "Plus tard : confirmer la signature du chèque (sans scan obligatoire dans COSUD) et renvoyer vers l’AC.",
                ],
            ),
            (
                "Autres courriers",
                [
                    "Circuit général : instruire, confier, orienter, valider une réponse selon l’étape.",
                    "Utiliser orientation / ventilation lorsque le dossier doit être dirigé vers une structure.",
                    "Surveiller la cloche pour les dossiers en attente DG.",
                ],
            ),
            (
                "Dettes & moratoires",
                [
                    "Consulter le suivi des factures / dettes avant d’instruire un moratoire.",
                    "Peut créer ou faire créer un plan de moratoire (permission create).",
                    "Les instructions écrites (scan) restent la référence pour le responsable suivi des dépenses.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Une notification SMS/cloche à l’enregistrement facture peut être coupée par l’admin (paramètre).",
                    "Le particulier DG aide au suivi mais n’exécute pas vos étapes de validation.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-directeur.pptx",
        "label": "Directeur",
        "mission": "Orienter, ventiler, signer ou rejeter les courriers relevant de sa direction.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Documents, Dossiers (y compris création structure / partage direction), Courriers, Organigramme, Tableau de bord.",
                    "Permissions courrier : orienter, ventiler, signer, rejeter, archiver — pas la création secrétariat classique.",
                    "Vue hiérarchique des documents de sa direction.",
                ],
            ),
            (
                "Actions courrier",
                [
                    "Lorsque vous êtes destinataire / acteur d’une étape : ouvrir la fiche et traiter (instruire, valider).",
                    "Orienter / ventiler vers le bon service ou agent.",
                    "Courrier départ : signer ou rejeter selon le circuit de validation.",
                    "Archiver lorsque le dossier est clos côté direction.",
                ],
            ),
            (
                "Documents & dossiers",
                [
                    "Créer et classer les documents de la direction.",
                    "Partager un dossier en lecture / écriture selon les besoins (partage direction).",
                    "Respecter la confidentialité : les dossiers confidentiels ont des règles d’accès strictes.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Le Bon pour accord facture prestataire est une étape DG, pas « directeur » de structure.",
                    "Vérifier l’étape affichée avant d’agir (le bouton disponible dépend de votre rôle et de l’étape).",
                ],
            ),
        ],
    },
    {
        "file": "Guide-admin.pptx",
        "label": "Administrateur",
        "mission": "Configurer COSUD, les droits, les référentiels et les paramètres système.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Accès complet : tous les menus métier + Administration (utilisateurs, paramètres, audit, rôles…).",
                    "Paramètres : structures, rôles, plan de classement, types, circuits courriers, workflow, catégories dépense.",
                    "Politique d’accès documents et Notifications courriers (interrupteurs métier).",
                ],
            ),
            (
                "Utilisateurs & rôles",
                [
                    "Créer les comptes, affecter structure(s) et rôle(s) Spatie.",
                    "Pour pouvoir créer un courrier arrivée : permission courriers.create + rattachement secrétariat (SEC).",
                    "Attribuer voir-factures / voir-depenses selon le cloisonnement souhaité.",
                    "Ne jamais partager le compte admin pour le travail quotidien métier.",
                ],
            ),
            (
                "Notifications & paramètres utiles",
                [
                    "Paramètres → Notifications courriers : activer / désactiver la notif DG (SMS + cloche) à l’enregistrement facture/MAD.",
                    "Paramètres → Politique d’accès documents : lecture dossier parent lors d’un envoi en validation.",
                    "Vérifier la configuration SMS (environnement) avec l’équipe technique.",
                ],
            ),
            (
                "Circuits & referentiels",
                [
                    "Circuits courriers : ne modifier une étape qu’avec impact métier validé.",
                    "Types de courriers, catégories de dépense, plan de classement : tenir à jour.",
                    "Audit : consulter les actions sensibles en cas d’incident.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Un changement de rôle / permission peut masquer des menus immédiatement.",
                    "Les migrations et seeders de rôles doivent rester alignés avec la production.",
                    "Sauvegardes et droits fichiers (scans) = responsabilité technique associée.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-utilisateur.pptx",
        "label": "Utilisateur",
        "mission": "Gérer ses documents et dossiers dans la GED, sans administration des courriers.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Menus : Tableau de bord, Documents, Dossiers.",
                    "Créer / modifier documents et dossiers dans le périmètre autorisé.",
                    "Pas de module Courriers ni de menus factures / MAD / moratoires.",
                ],
            ),
            (
                "Travail quotidien",
                [
                    "Déposer et classer les documents dans le bon dossier.",
                    "Respecter les règles de confidentialité et de partage.",
                    "Participer aux validations documentaires si vous êtes désigné dans un circuit.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Si vous devez traiter un courrier, un rôle secrétariat / circuit doit vous être attribué.",
                    "La suppression de dossiers / documents peut être limitée selon les droits.",
                ],
            ),
        ],
    },
    {
        "file": "Guide-chef-service.pptx",
        "label": "Chef de service",
        "mission": "Piloter documents et dossiers de son service, avec accès à l’organigramme.",
        "slides": [
            (
                "Périmètre & menus",
                [
                    "Comme Utilisateur : Documents, Dossiers, Tableau de bord.",
                    "En plus : Organigramme (consultation de la structure).",
                    "Pas de droits courriers métier dans le profil de base.",
                ],
            ),
            (
                "Travail quotidien",
                [
                    "Organiser le classement du service (dossiers, documents).",
                    "Valider ou transmettre les documents selon les circuits GED.",
                    "S’appuyer sur l’organigramme pour identifier les interlocuteurs.",
                ],
            ),
            (
                "Points d’attention",
                [
                    "Les rôles chef_projet / chef_pool / chef_centre ont le même socle documentaire (sans organigramme).",
                    "Pour le courrier, demander l’attribution d’un rôle circuit adapté.",
                ],
            ),
        ],
    },
]


def build_guide(guide: dict) -> Path:
    prs = new_prs()
    add_title_slide(prs, guide["label"], guide["mission"])
    for title, bullets in guide["slides"]:
        add_section(prs, guide["label"], title, bullets)
    # slide de clôture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_bg(slide, DARK)
    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.5), Inches(2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Fin du guide"
    _set_run(r, size=32, bold=True, color=WHITE)
    p2 = tf.add_paragraph()
    r2 = p2.add_run()
    r2.text = f"{guide['label']} — COSUD / ACSI"
    _set_run(r2, size=18, color=GREEN)
    p2.space_before = Pt(16)
    p3 = tf.add_paragraph()
    r3 = p3.add_run()
    r3.text = "Pour toute évolution de droits, s’adresser à l’administrateur COSUD."
    _set_run(r3, size=14, italic=True, color=LIGHT)
    p3.space_before = Pt(20)

    OUT_DIR.mkdir(parents=True, exist_ok=True)
    path = OUT_DIR / guide["file"]
    prs.save(str(path))
    return path


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    for guide in GUIDES:
        path = build_guide(guide)
        print(f"OK: {path}")
    print(f"Total: {len(GUIDES)} guides -> {OUT_DIR}")


if __name__ == "__main__":
    main()
