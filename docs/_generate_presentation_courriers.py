"""Génère docs/Presentation-module-gestion-courriers.pptx — présentation direction + schéma architecture."""

from __future__ import annotations

import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

GREEN = RGBColor(0x00, 0xA0, 0x55)
GREEN_DARK = RGBColor(0x06, 0x7A, 0x45)
DARK = RGBColor(0x0F, 0x17, 0x2A)
SLATE = RGBColor(0x33, 0x41, 0x55)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF1, 0xF5, 0xF9)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SOFT = RGBColor(0xCB, 0xD5, 0xE1)
BLUE_SOFT = RGBColor(0xE0, 0xF2, 0xFE)
GREEN_SOFT = RGBColor(0xDC, 0xFC, 0xE7)
AMBER_SOFT = RGBColor(0xFF, 0xF7, 0xED)

DOCS_DIR = os.path.dirname(__file__)
LOGO_PATH = os.path.normpath(os.path.join(DOCS_DIR, "..", "public", "images", "image-logo.jpg"))


def add_bg(slide, color=WHITE) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    sp_tree = slide.shapes._spTree
    sp = shape._element
    sp_tree.remove(sp)
    sp_tree.insert(2, sp)


def set_run(paragraph, text: str, size: int = 18, bold: bool = False, color=SLATE) -> None:
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"


def add_header_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.9))
    bar.fill.solid()
    bar.fill.fore_color.rgb = GREEN
    bar.line.fill.background()

    foot = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(7.15), prs.slide_width, Inches(0.35)
    )
    foot.fill.solid()
    foot.fill.fore_color.rgb = DARK
    foot.line.fill.background()
    p = foot.text_frame.paragraphs[0]
    set_run(
        p,
        "GED ACSI  |  Présentation Direction  |  Digitalisation des courriers DG",
        10,
        False,
        WHITE,
    )
    p.alignment = PP_ALIGN.CENTER


def add_title(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12.3), Inches(0.55))
    set_run(box.text_frame.paragraphs[0], text, 26, True, WHITE)


def add_bullets(slide, items: list[str], left=0.5, top=1.15, width=12.3, height=5.7, size=16) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(size)
        p.font.color.rgb = SLATE
        p.font.name = "Calibri"
        p.space_after = Pt(6)


def add_dual_column(slide, left_title, left_items, right_title, right_items, size=14) -> None:
    card1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.15), Inches(6.1), Inches(5.6)
    )
    card1.fill.solid()
    card1.fill.fore_color.rgb = LIGHT
    card1.line.fill.background()
    t1 = slide.shapes.add_textbox(Inches(0.65), Inches(1.3), Inches(5.6), Inches(0.4))
    set_run(t1.text_frame.paragraphs[0], left_title, 18, True, GREEN)
    add_bullets(slide, left_items, 0.65, 1.8, 5.6, 4.7, size)

    card2 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.15), Inches(6.1), Inches(5.6)
    )
    card2.fill.solid()
    card2.fill.fore_color.rgb = LIGHT
    card2.line.fill.background()
    t2 = slide.shapes.add_textbox(Inches(7.05), Inches(1.3), Inches(5.6), Inches(0.4))
    set_run(t2.text_frame.paragraphs[0], right_title, 18, True, GREEN)
    add_bullets(slide, right_items, 7.05, 1.8, 5.6, 4.7, size)


def content_slide(title: str, bullets: list[str], size: int = 16):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide)
    add_title(slide, title)
    add_bullets(slide, bullets, size=size)
    return slide


def add_diagram_box(
    slide,
    left: float,
    top: float,
    width: float,
    height: float,
    title: str,
    subtitle: str = "",
    fill=LIGHT,
    title_color=GREEN_DARK,
    title_size: int = 12,
    subtitle_size: int = 10,
) -> None:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = GREEN
    shape.line.width = Pt(1.25)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, title, title_size, True, title_color)
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2, subtitle, subtitle_size, False, SLATE)


def add_diagram_arrow(slide, x1: float, y1: float, x2: float, y2: float) -> None:
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(x1),
        Inches(y1),
        Inches(x2),
        Inches(y2),
    )
    conn.line.color.rgb = GREEN_DARK
    conn.line.width = Pt(2.5)
    conn.line.end_arrowhead = True


def add_architecture_diagram_slide() -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide)
    add_title(slide, "4. Architecture de la solution — vue d'ensemble")

    # Bandeau utilisateurs
    users = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.05), Inches(12.3), Inches(0.85)
    )
    users.fill.solid()
    users.fill.fore_color.rgb = DARK
    users.line.fill.background()
    tf = users.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "Utilisateurs — DG, Secrétariat, Directions, Agence comptable, Partenaires internes", 13, True, WHITE)

    add_diagram_arrow(slide, 6.65, 1.9, 6.65, 2.15)

    # Plateforme GED
    platform = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(2.2), Inches(12.3), Inches(2.55)
    )
    platform.fill.solid()
    platform.fill.fore_color.rgb = GREEN_SOFT
    platform.line.color.rgb = GREEN
    platform.line.width = Pt(2)
    tf = platform.text_frame
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.08)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(p, "PLATEFORME GED ACSI — accès web sécurisé", 14, True, GREEN_DARK)

    modules = [
        (0.75, 2.75, 2.85, 0.95, "Documents", "Archivage & partage"),
        (3.75, 2.75, 2.85, 0.95, "Dossiers", "Classement & métadonnées"),
        (6.75, 2.75, 2.85, 0.95, "Recherche", "Retrouver rapidement"),
        (9.75, 2.75, 2.85, 0.95, "Administration", "Utilisateurs & paramètres"),
    ]
    for left, top, w, h, title, sub in modules:
        add_diagram_box(slide, left, top, w, h, title, sub, WHITE, GREEN_DARK, 11, 9)

    # Module courriers (mis en avant)
    courrier_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(3.95), Inches(11.3), Inches(0.65)
    )
    courrier_box.fill.solid()
    courrier_box.fill.fore_color.rgb = GREEN
    courrier_box.line.fill.background()
    tf = courrier_box.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p,
        "MODULE GESTION DES COURRIERS — Registres  |  Circuits métier  |  Suivi paiements (FSP)  |  Notifications",
        12,
        True,
        WHITE,
    )

    add_diagram_arrow(slide, 6.65, 4.75, 6.65, 5.0)

    # Couche données
    data = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(5.05), Inches(12.3), Inches(0.95)
    )
    data.fill.solid()
    data.fill.fore_color.rgb = BLUE_SOFT
    data.line.color.rgb = GREEN
    data.line.width = Pt(1.5)
    tf = data.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    set_run(
        p,
        "Données & services — Base sécurisée  |  Fichiers numérisés  |  E-mail & alertes  |  Journal de traçabilité",
        12,
        True,
        SLATE,
    )

    # Légende flux courrier (bas)
    legend = slide.shapes.add_textbox(Inches(0.55), Inches(6.15), Inches(12.2), Inches(0.85))
    tf = legend.text_frame
    tf.word_wrap = True
    set_run(
        tf.paragraphs[0],
        "Flux courrier : Arrivée enregistrée → circuit métier (facture ou courrier général) → traitement → Départ signé → archivage",
        11,
        False,
        MUTED,
    )


def add_courrier_flow_diagram_slide() -> None:
    """Schéma simplifié du parcours d'un courrier — lisible par la direction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide)
    add_title(slide, "7. Parcours d'un courrier — vue simplifiée")

    steps_a = [
        ("Enregistrement\nArrivée", 0.55),
        ("Instruction\nDG", 2.55),
        ("Traitement\ndossier / AC", 4.55),
        ("Chèque &\nFSP", 6.55),
        ("Clôture\npaiement", 8.55),
    ]
    steps_b = [
        ("Enregistrement\nArrivée", 0.55),
        ("Instruction\nDG", 2.55),
        ("Projet de\nréponse", 4.55),
        ("Validation\n& Départ", 6.55),
        ("Expédition\n& archivage", 8.55),
    ]

    # Labels circuits
    la = slide.shapes.add_textbox(Inches(0.55), Inches(1.05), Inches(5.5), Inches(0.35))
    set_run(la.text_frame.paragraphs[0], "Circuit A — Factures & MAD", 13, True, GREEN_DARK)
    lb = slide.shapes.add_textbox(Inches(0.55), Inches(3.55), Inches(5.5), Inches(0.35))
    set_run(lb.text_frame.paragraphs[0], "Circuit B — Courriers généraux", 13, True, GREEN_DARK)

    def draw_flow(steps: list[tuple[str, float]], top: float, fill) -> None:
        prev_x = None
        for i, (label, left) in enumerate(steps):
            add_diagram_box(slide, left, top, 1.65, 0.95, label.split("\n")[0], label.split("\n")[1] if "\n" in label else "", fill, GREEN_DARK, 10, 9)
            if prev_x is not None:
                add_diagram_arrow(slide, prev_x + 1.65, top + 0.48, left, top + 0.48)
            prev_x = left

    draw_flow(steps_a, 1.45, AMBER_SOFT)
    draw_flow(steps_b, 3.95, GREEN_SOFT)

    # Encadré bénéfices
    note = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(5.35), Inches(12.2), Inches(1.55)
    )
    note.fill.solid()
    note.fill.fore_color.rgb = LIGHT
    note.line.fill.background()
    tf = note.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.15)
    set_run(tf.paragraphs[0], "À chaque étape : responsable identifié, délai suivi, notification automatique, historique conservé.", 13, True, GREEN_DARK)
    p2 = tf.add_paragraph()
    set_run(p2, "Le DG et le secrétariat disposent d'une vision en temps réel de l'avancement de chaque dossier.", 12, False, SLATE)


def add_roadmap_slide() -> None:
    """Feuille de route visuelle — 4 phases pour la Direction."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_header_bar(slide)
    add_title(slide, "15. Feuille de route — calendrier proposé")

    phases = [
        (
            "Phase 1",
            "Réalisée",
            "Juillet 2026",
            [
                "Module courriers livré",
                "Circuits A & B opérationnels",
                "FSP & registres numériques",
                "Tests validés (env. test)",
            ],
            GREEN_SOFT,
            True,
        ),
        (
            "Phase 2",
            "Immédiate",
            "Août 2026",
            [
                "Formation des acteurs DG",
                "Mise en production DG",
                "Rodage & support SI",
                "Registres papier en parallèle",
            ],
            AMBER_SOFT,
            False,
        ),
        (
            "Phase 3",
            "Court terme",
            "T4 2026",
            [
                "Bilan d'usage Direction",
                "Ajustements circuits / délais",
                "Tableaux de bord délais",
                "Extension 1–2 directions",
            ],
            BLUE_SOFT,
            False,
        ),
        (
            "Phase 4",
            "Moyen terme",
            "2027",
            [
                "Généralisation ACSI",
                "Signature électronique",
                "Indicateurs de pilotage",
                "Archivage pérenne",
            ],
            LIGHT,
            False,
        ),
    ]

    card_w = 2.9
    gap = 0.25
    start_x = 0.55
    top = 1.2

    for i, (phase, statut, periode, items, fill, done) in enumerate(phases):
        left = start_x + i * (card_w + gap)
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(5.5)
        )
        card.fill.solid()
        card.fill.fore_color.rgb = fill
        card.line.color.rgb = GREEN if done else SOFT
        card.line.width = Pt(2 if done else 1)

        # bandeau phase
        band = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(card_w), Inches(0.85)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = GREEN if done else DARK
        band.line.fill.background()
        tf = band.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, phase, 14, True, WHITE)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        set_run(p2, statut, 11, False, SOFT)

        # période
        per = slide.shapes.add_textbox(Inches(left + 0.1), Inches(top + 1.0), Inches(card_w - 0.2), Inches(0.35))
        p = per.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        set_run(p, periode, 13, True, GREEN_DARK)

        add_bullets(slide, [f"• {it}" for it in items], left + 0.15, top + 1.5, card_w - 0.3, 3.8, 12)

        if i < len(phases) - 1:
            add_diagram_arrow(slide, left + card_w, top + 2.7, left + card_w + gap, top + 2.7)


# ===== COVER =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
accent.fill.solid()
accent.fill.fore_color.rgb = GREEN
accent.line.fill.background()
if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Inches(0.9), Inches(0.55), height=Inches(0.85))
box = s.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.5), Inches(0.5))
set_run(box.text_frame.paragraphs[0], "ACSI — Direction Générale", 22, True, GREEN)
box2 = s.shapes.add_textbox(Inches(0.9), Inches(2.35), Inches(11.5), Inches(1.3))
set_run(box2.text_frame.paragraphs[0], "Digitalisation de la gestion des courriers", 38, True, WHITE)
box3 = s.shapes.add_textbox(Inches(0.9), Inches(3.7), Inches(11.5), Inches(1.8))
tf = box3.text_frame
tf.word_wrap = True
set_run(tf.paragraphs[0], "Projet GED — Module Gestion des courriers", 22, False, SOFT)
p2 = tf.add_paragraph()
set_run(p2, "Présentation à la Direction — Contexte, solution et bénéfices", 16, False, MUTED)
p3 = tf.add_paragraph()
set_run(p3, "Juillet 2026", 14, False, MUTED)

# ===== 1 CONTEXTE =====
content_slide(
    "1. Pourquoi ce projet ?",
    [
        "Le secrétariat de la DG tient aujourd'hui des registres papier Arrivée / Départ",
        "Les circuits de traitement (factures, courriers, chèques) reposent sur des échanges informels",
        "Conséquences : délais difficiles à suivre, risque de perte d'information, faible visibilité pour la Direction",
        "La digitalisation vise à conserver le processus métier existant, en le rendant traçable et pilotable",
        "Le projet s'inscrit dans la modernisation globale de la GED ACSI à la Direction Générale",
    ],
    17,
)

# ===== 2 OBJECTIFS =====
content_slide(
    "2. Objectifs pour la Direction",
    [
        "Savoir à tout moment où en est chaque courrier et qui doit agir",
        "Réduire les retards grâce aux alertes et relances automatiques",
        "Conserver l'équivalent numérique des registres officiels Arrivée / Départ",
        "Sécuriser le circuit des factures : du dossier au chèque, puis à la clôture du paiement",
        "Faciliter le travail du secrétariat, du DG et de l'agence comptable sans changer leurs rôles",
        "Préparer une extension progressive à d'autres services de l'organisme",
    ],
    17,
)

# ===== 3 VISION GED =====
content_slide(
    "3. La GED ACSI — une plateforme unique",
    [
        "Application web accessible aux agents habilités de la DG et services associés",
        "Centralise documents, dossiers, recherche et administration des accès",
        "Le module « Gestion des courriers » est le premier volet métier prioritaire livré",
        "Même identité visuelle, mêmes règles de sécurité, même organigramme pour tous les modules",
        "Hébergement sur serveur interne ACSI — données maîtrisées par l'organisme",
        "Application distincte de SIFEC (état civil) — périmètre et utilisateurs différents",
    ],
    17,
)

# ===== 4 ARCHITECTURE VISUELLE =====
add_architecture_diagram_slide()

# ===== 5 MODULE COURRIERS =====
content_slide(
    "5. Module Gestion des courriers — ce qu'il apporte",
    [
        "Registre numérique Arrivée : numérotation, expéditeur, objet, lien vers la réponse",
        "Registre numérique Départ : destinataire, date, pièces jointes, statut d'expédition",
        "Circuits métier paramétrables : deux parcours natifs (factures et courriers généraux)",
        "Suivi des paiements (FSP) : tableau automatique des factures et mises à disposition",
        "Notifications par cloche et e-mail à chaque étape importante",
        "Historique complet : qui a fait quoi, quand, avec quelles instructions",
    ],
    16,
)

# ===== 6 ACTEURS =====
content_slide(
    "6. Acteurs et responsabilités",
    [
        "Secrétariat DG — Enregistre les arrivées, pilote les registres, crée les départs",
        "Directeur Général — Donne les instructions, valide, signe (courriers et chèques)",
        "Particulière du DG — Traite les dossiers, prépare les projets de réponse",
        "Responsable dossiers prestataires — Traite les factures avant l'agence comptable",
        "Agence comptable (AC) — Établit le chèque, gère caissiers et retour de caisse",
        "Responsable suivi des dépenses — Suit les FSP et clôture après preuve de paiement",
        "Directeurs de direction — Peuvent instruire si leur structure est destinataire",
    ],
    15,
)

# ===== 7 PARCOURS VISUEL =====
add_courrier_flow_diagram_slide()

# ===== 8 CIRCUITS DETAIL =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_header_bar(s)
add_title(s, "8. Les deux circuits métier en détail")
add_dual_column(
    s,
    "Circuit A — Factures & MAD",
    [
        "Factures prestataires et mises à disposition (MAD)",
        "Étape 1 : enregistrement → instruction DG → traitement dossier",
        "Étape 2 : l'AC établit le chèque (montant obligatoire) → signature DG",
        "Création automatique de la ligne FSP + alerte responsable dépenses",
        "Relais automatiques vers caissiers et retour de caisse",
        "Clôture après dépôt de la preuve de paiement",
    ],
    "Circuit B — Courriers généraux",
    [
        "Courriers administratifs, invitations, demandes, notes…",
        "Enregistrement → instruction DG → projet de réponse (particulière)",
        "Validation DG → création du courrier de départ",
        "Signature, expédition, réception par le secrétariat destinataire",
        "Variantes : réponse directe du DG, rejet du projet, instruction par directeur",
        "Archivage une fois le cycle terminé",
    ],
    14,
)

# ===== 9 FSP =====
content_slide(
    "9. Suivi des paiements (FSP) — pilotage financier",
    [
        "Écran dédié « Suivi des paiements » dans le menu GED",
        "Deux tableaux : FSP FACTURE (prestataires) et FSP MAD (mises à disposition)",
        "Alimentation automatique dès l'envoi du chèque au DG (aucune ressaisie)",
        "Colonnes : numéro, date, objet, montant, fournisseur, instructions du DG",
        "Export Excel pour contrôle, reporting et archivage",
        "Responsable suivi des dépenses notifiée à l'entrée chèque, clôture à la preuve paiement",
    ],
    16,
)

# ===== 10 REGISTRES =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s)
add_header_bar(s)
add_title(s, "10. Registres numériques — continuité avec le papier")
add_dual_column(
    s,
    "Registre Arrivée",
    [
        "Reproduit les colonnes du registre papier officiel",
        "Date et numéro d'arrivée",
        "Expéditeur et objet du courrier",
        "Lien automatique vers le courrier de réponse (départ)",
        "Impression et export pour archivage",
    ],
    "Registre Départ",
    [
        "Numéro d'ordre et date de départ",
        "Destinataire, objet, nombre de pièces",
        "Numéro d'archives et observations",
        "Statut : brouillon, signé, expédié, reçu",
        "Traçabilité de bout en bout",
    ],
    15,
)

# ===== 11 NOTIFICATIONS =====
content_slide(
    "11. Alertes et suivi des délais",
    [
        "Chaque acteur est alerté lorsqu'un courrier arrive à son étape",
        "Cloche dans l'application + e-mail (expéditeur « GED ACSI »)",
        "Alerte retard si une étape dépasse le délai configuré (48 h par défaut)",
        "Possibilité de relancer le responsable en retard (hors tour du DG)",
        "Entrée chèque : notification immédiate au responsable suivi des dépenses",
        "Courriers confidentiels : diffusion restreinte aux seuls acteurs concernés",
    ],
    16,
)

# ===== 12 SECURITE =====
content_slide(
    "12. Sécurité et gouvernance",
    [
        "Accès par compte personnel — chaque agent ne voit que ce qui le concerne",
        "Rôles métier : DG, secrétariat, AC, responsable dépenses, directeurs…",
        "Courriers confidentiels : visibilité limitée, pas de notification généralisée",
        "Journal d'audit : actions sensibles enregistrées pour contrôle ultérieur",
        "Authentification renforcée disponible (double facteur) pour les comptes sensibles",
        "Données hébergées en interne — pas de cloud externe",
    ],
    16,
)

# ===== 13 BENEFICES =====
content_slide(
    "13. Bénéfices concrets pour la Direction",
    [
        "Visibilité en temps réel : plus de « où en est ce dossier ? » sans réponse",
        "Responsabilisation : chaque étape a un acteur identifié et daté",
        "Gain de temps secrétariat : moins de relances manuelles, registres auto-alimentés",
        "Maîtrise des dépenses : FSP synchronisé avec le circuit facture/chèque",
        "Conformité : registres numériques alignés sur les registres papier existants",
        "Évolutivité : circuits modifiables sans refonte complète du système",
    ],
    16,
)

# ===== 14 DEPLOIEMENT =====
content_slide(
    "14. Mise en service et accompagnement",
    [
        "Phase actuelle : module courriers opérationnel pour la DG (environnement de test validé)",
        "Comptes utilisateurs créés selon l'organigramme et les rôles métier",
        "Formation prévue : secrétariat, DG, agence comptable, responsable dépenses",
        "Documentation disponible : scénarios de test, résumé des processus, présentation",
        "Accompagnement au démarrage : période de rodage avec support SI",
        "Extension future : autres directions, signature électronique, tableaux de bord",
    ],
    16,
)

# ===== 15 FEUILLE DE ROUTE =====
add_roadmap_slide()

# ===== CLOSING =====
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, DARK)
accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.25), prs.slide_height)
accent.fill.solid()
accent.fill.fore_color.rgb = GREEN
accent.line.fill.background()
if os.path.exists(LOGO_PATH):
    s.shapes.add_picture(LOGO_PATH, Inches(0.9), Inches(0.7), height=Inches(0.75))
box = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(0.6))
set_run(box.text_frame.paragraphs[0], "En résumé", 22, True, GREEN)
box2 = s.shapes.add_textbox(Inches(0.9), Inches(2.7), Inches(11.5), Inches(2.8))
tf = box2.text_frame
tf.word_wrap = True
set_run(tf.paragraphs[0], "Une solution GED intégrée pour moderniser la DG", 28, True, WHITE)
p2 = tf.add_paragraph()
set_run(p2, "Registres numériques · Circuits métier · Suivi FSP · Alertes automatiques", 18, False, SOFT)
p3 = tf.add_paragraph()
set_run(p3, "Le processus métier est préservé — la traçabilité et le pilotage sont renforcés", 18, False, SOFT)
p4 = tf.add_paragraph()
set_run(p4, "", 12, False, WHITE)
p5 = tf.add_paragraph()
set_run(p5, "Questions · Démonstration · Validation de la feuille de route", 16, False, MUTED)

out = os.path.join(DOCS_DIR, "Presentation-module-gestion-courriers.pptx")
tmp = os.path.join(DOCS_DIR, "_Presentation-module-gestion-courriers.tmp.pptx")
prs.save(tmp)
try:
    if os.path.exists(out):
        os.remove(out)
    os.replace(tmp, out)
except OSError as e:
    fallback = os.path.join(DOCS_DIR, "Presentation-module-gestion-courriers-NEW.pptx")
    os.replace(tmp, fallback)
    print(f"ATTENTION: fichier cible verrouille ({e}). Sauvegarde: {fallback}")
    raise SystemExit(0)
print(f"OK {out} ({len(prs.slides)} slides)")
print(f"Logo: {'oui' if os.path.exists(LOGO_PATH) else 'introuvable — ' + LOGO_PATH}")
